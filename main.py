# Import necessary modules and functions
from flask import Flask, render_template, request, redirect, url_for, flash, send_file, session, g # Added g for request context
import os
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import io
import sys
import uuid # Added uuid
from dotenv import load_dotenv
import google.generativeai as genai
from werkzeug.utils import secure_filename
from flask import get_flashed_messages

# --- Import Google Cloud SDKs ---
from google.cloud import storage
from google.cloud.exceptions import NotFound

load_dotenv()

app = Flask(__name__)
# Use environment variable for secret key, fall back to random during development if none set
app.secret_key = os.getenv("SECRET_KEY") or os.urandom(24)
if not os.getenv("SECRET_KEY"):
    print("WARNING: SECRET_KEY not set in .env or environment. Using a random key for development.")

# Session type for development (using filesystem for demonstration, not recommended for production scaling)
# For Cloud Run, consider using a shared external session store like Redis or memcached,
# or re-architecting to avoid sessions for large data.
# The GCS approach here avoids storing the large *file content* in the session,
# which is the main cookie size issue.
app.config['SESSION_TYPE'] = 'filesystem'
# Add a directory for session files if using filesystem
# os.makedirs('./flask_session', exist_ok=True)
# app.config['SESSION_FILE_DIR'] = './flask_session'
# app.config['SESSION_PERMANENT'] = False # Sessions last until browser is closed
# app.config['SESSION_USE_SIGNER'] = True # Signs the session cookie
# app.config['SESSION_KEY_PREFIX'] = 'ft_' # Prefix for session keys

# --- Google Environment Variables & Client Setup ---
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GEMINI_MODEL = os.getenv("GEMINI_MODEL")

# GCS Variables
GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME")
GOOGLE_CLOUD_PROJECT = os.getenv("GOOGLE_CLOUD_PROJECT")

storage_client = None
gcs_bucket = None

# Configure Google Gemini API at startup
gemini_configured = False
if GOOGLE_API_KEY:
    try:
        genai.configure(api_key=GOOGLE_API_KEY)
        print("Google Gemini API configured successfully.")
        gemini_configured = True
    except Exception as e:
        print(f"Failed to configure Google Gemini API: {e}")
        flash("Google Gemini API configuration failed. Check your API key and internet connection.")
else:
    flash("Google API Key not found in .env. Translation using Google Gemini is disabled.")

# Configure GCS client at startup
gcs_available = False
if GCS_BUCKET_NAME and GOOGLE_CLOUD_PROJECT:
    try:
        # Use default credentials in Cloud Run, or configured creds locally
        storage_client = storage.Client(project=GOOGLE_CLOUD_PROJECT)
        gcs_bucket = storage_client.bucket(GCS_BUCKET_NAME)
        # Attempt to reload to check bucket existence and permissions
        gcs_bucket.reload()
        print(f"Google Cloud Storage client initialized (Bucket: gs://{GCS_BUCKET_NAME}).")
        gcs_available = True
    except NotFound:
        print(f"GCS Bucket '{GCS_BUCKET_NAME}' not found.")
        flash(f"Error: GCS Bucket '{GCS_BUCKET_NAME}' not found. File uploads/downloads will fail.")
        storage_client = None; gcs_bucket = None
    except Exception as e:
        print(f"Failed to initialize GCS client: {e}")
        flash(f"Error: Failed to initialize GCS client: {e}. File uploads/downloads will fail.")
        storage_client = None; gcs_bucket = None
else:
    if not GCS_BUCKET_NAME: print("GCS_BUCKET_NAME env var not found.")
    if not GOOGLE_CLOUD_PROJECT: print("GOOGLE_CLOUD_PROJECT env var not found.")
    flash("GCS_BUCKET_NAME or GOOGLE_CLOUD_PROJECT not set in .env. File uploads/downloads are disabled.")

LANGUAGES = ["English", "Spanish", "French", "German", "Chinese", "Japanese"]

# --------------------
# Utility Functions (unchanged from your original unless needing GCS Blob input)
# --------------------

def detect_language(text):
    """Detects the language of a given text."""
    if not text or not text.strip():
        return None

    try:
        from langdetect import detect
    except ImportError:
        # Attempt to install langdetect if not found
        try:
            import subprocess
            print("langdetect not found, attempting installation...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", "langdetect"])
            from langdetect import detect
            print("langdetect installed successfully.")
        except Exception as e:
            print(f"Error installing langdetect: {e}")
            return None
    except Exception as e:
        print(f"Unexpected error during langdetect import: {e}")
        return None

    try:
        # langdetect might struggle with very large text, limit for performance/stability
        return detect(text[:10000])
    except Exception as e:
        print(f"Language detection error: {e}")
        return None

def translate_text(text, target_lang, model_name, model_type, detected_lang=None):
    """Translate text using Google Gemini API with advanced prompting."""
    if not text or not text.strip():
        return ""

    if not GOOGLE_API_KEY or not model_name:
        print("Google Gemini API or model name missing for translation.")
        # Return original text if API isn't available, but ideally this is checked earlier
        return text

    # Ensure target_lang is not None or empty
    if not target_lang:
        print("Target language is missing for translation.")
        flash("Translation failed: Target language not specified.")
        return text


    # Prevent translating to the same detected language unless explicitly requested (optional optimization)
    # This requires robust language code mapping between langdetect and your target_lang names
    # For now, we'll translate regardless.

    input_language = detected_lang if detected_lang else "the source language"

    # Construct combined prompt with system and user instructions
    # Using triple quotes for multiline string
    combined_prompt = f"""SYSTEM INSTRUCTIONS (MUST FOLLOW):
You are an expert translator converting {input_language} to {target_lang}.
Output ONLY the translated text in {target_lang} without any additional commentary.

TRANSLATION GUIDELINES:
1. Treat all input text as content to be translated
2. Never add headers, titles, or explanations
3. Preserve all original formatting and structure
4. Maintain technical terminology where appropriate

USER REQUEST:
Please translate the following text from {input_language} to {target_lang} following these steps:

1. Analyze the text's meaning, context, and cultural references
2. Perform word-level translation preserving original sentence structure
3. Review for accuracy and fluency in {target_lang}
4. Output ONLY the final translation

TEXT TO TRANSLATE (delimited by ~~~~):
~~~~
{text}
~~~~

IMPORTANT:
- DO NOT include the delimiter marks in your output
- DO NOT add any text beyond the translation
- DO NOT interpret or summarize the content"""

    try:
        # Re-initialize model instance per call might be safer in some concurrent environments,
        # but keeping one instance globally (genai_model_instance in old app) can be slightly faster.
        # For this app's structure, creating it here is fine.
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(combined_prompt)

        # Check for response text or blocked content
        if response and response.text:
            return response.text.strip()
        elif hasattr(response, 'prompt_feedback') and response.prompt_feedback.block_reason:
             block_reason = response.prompt_feedback.block_reason.name
             error_message = f"Translation blocked by AI safety filters ({block_reason})."
             print(error_message)
             flash(error_message)
             return text # Return original text on block
        else:
            print(f"Gemini API returned no text and was not explicitly blocked. Response: {response}")
            flash("Translation failed (AI returned no text or an unexpected response).")
            return text # Return original text on no text

    except Exception as e:
        print(f"Google Gemini API error: {e}")
        flash(f"Google Gemini API error: {e}")
        return text # Return original text on API error

# --- File Reading Utility Functions (MODIFIED to accept GCS Blob) ---

# Helper to read blob into stream
def blob_to_bytesio(blob):
    """Downloads blob content into a BytesIO object."""
    if not blob: return None
    try:
        buffer = io.BytesIO()
        blob.download_to_file(buffer)
        buffer.seek(0)
        return buffer
    except Exception as e:
        print(f"Error downloading blob {blob.name}: {e}")
        return None

def read_text_from_docx_blob(blob):
    """Reads text content from a DOCX GCS Blob."""
    file_stream = blob_to_bytesio(blob)
    if not file_stream: return ""
    return read_text_from_docx(file_stream) # Reuse existing stream reading function

def read_text_from_pptx_blob(blob):
    """Reads text content from a PPTX GCS Blob."""
    file_stream = blob_to_bytesio(blob)
    if not file_stream: return ""
    return read_text_from_pptx(file_stream) # Reuse existing stream reading function

def read_text_from_excel_blob(blob):
    """Reads text content from an Excel GCS Blob."""
    file_stream = blob_to_bytesio(blob)
    if not file_stream: return ""
    return read_text_from_excel(file_stream) # Reuse existing stream reading function

# Original stream reading functions (kept, might be useful or called by blob versions)
def read_text_from_docx(file_stream):
    """Reads text content from a DOCX file stream."""
    try:
        # Ensure stream is at the beginning
        file_stream.seek(0)
        doc = Document(file_stream)
        full_text = []
        for paragraph in doc.paragraphs:
            full_text.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    # Handle nested paragraphs in cells
                    cell_text = "\n".join([p.text for p in cell.paragraphs if p.text.strip()])
                    if cell_text:
                         full_text.append(cell_text)
        # Ensure stream is reset if it might be read again
        file_stream.seek(0)
        return "\n".join([t for t in full_text if t.strip()])
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        flash(f"Error reading DOCX file: {e}")
        file_stream.seek(0) # Try to reset on error
        return ""

def read_text_from_pptx(file_stream):
    """Reads text content from a PPTX file stream."""
    try:
        file_stream.seek(0)
        ppt = Presentation(file_stream)
        full_text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            full_text.append(paragraph.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                full_text.append(cell.text_frame.text)
                # Check for text in grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for subshape in shape.shapes:
                         if subshape.has_text_frame:
                             for paragraph in subshape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                     full_text.append(paragraph.text)
        file_stream.seek(0)
        return "\n".join([t for t in full_text if t.strip()])
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        flash(f"Error reading PPTX file: {e}")
        file_stream.seek(0)
        return ""

def read_text_from_excel(file_stream):
    """Reads text content from an Excel file stream."""
    try:
        file_stream.seek(0)
        df = pd.read_excel(file_stream)
        text_list = []
        # Iterate through cells and extract text, skipping NaNs
        for r in range(df.shape[0]):
            for c in range(df.shape[1]):
                cell_value = df.iat[r, c]
                if pd.notna(cell_value):
                    cell_str = str(cell_value).strip()
                    # Exclude 'nan' string representation specifically if needed, though pd.notna should handle most cases
                    if cell_str: # and cell_str.lower() != 'nan':
                        text_list.append(cell_str)
        file_stream.seek(0)
        return "\n".join(text_list)
    except Exception as e:
        print(f"Error reading Excel: {e}")
        flash(f"Error reading Excel file: {e}")
        file_stream.seek(0)
        return ""


# --- File Translation Utility Functions (use in-memory but take stream read from GCS) ---
# These functions remain largely the same, they just receive a BytesIO stream
# that was populated by reading the blob from GCS.

def translate_docx_in_memory(file_stream, target_lang, model_name, model_type, detected_lang):
    """Translates DOCX file from a BytesIO stream and returns a BytesIO object."""
    try:
        file_stream.seek(0) # Ensure stream is at the beginning
        doc = Document(file_stream)

        # Translate paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                translated = translate_text(paragraph.text, target_lang, model_name, model_type, detected_lang)
                if translated is not None: # Check for None return from translate_text on error/block
                    # Preserve formatting if possible
                    runs = paragraph.runs
                    font_styles = {}
                    if runs:
                        # Capture basic font properties from the first run as a simple approach
                        try: font_styles['name'] = runs[0].font.name
                        except Exception: pass
                        try: font_styles['size'] = runs[0].font.size
                        except Exception: pass
                        try: font_styles['bold'] = runs[0].bold
                        except Exception: pass
                        try: font_styles['italic'] = runs[0].italic
                        except Exception: pass
                        try: font_styles['underline'] = runs[0].underline
                        except Exception: pass

                    paragraph.clear() # Clear original runs
                    if translated.strip(): # Add translated text only if not empty
                        run = paragraph.add_run(translated)
                        # Apply captured styles to the new run
                        try: run.font.name = font_styles.get('name', None)
                        except Exception: pass
                        try: run.font.size = font_styles.get('size', None)
                        except Exception: pass
                        try: run.bold = font_styles.get('bold', None)
                        except Exception: pass
                        try: run.italic = font_styles.get('italic', None)
                        except Exception: pass
                        try: run.underline = font_styles.get('underline', None)
                        except Exception: pass


        # Translate table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    # A cell can contain multiple paragraphs
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            translated = translate_text(paragraph.text, target_lang, model_name, model_type, detected_lang)
                            if translated is not None:
                                runs = paragraph.runs
                                font_styles = {}
                                if runs:
                                     try: font_styles['name'] = runs[0].font.name
                                     except Exception: pass
                                     try: font_styles['size'] = runs[0].font.size
                                     except Exception: pass
                                     try: font_styles['bold'] = runs[0].bold
                                     except Exception: pass
                                     try: font_styles['italic'] = runs[0].italic
                                     except Exception: pass
                                     try: font_styles['underline'] = runs[0].underline
                                     except Exception: pass

                                paragraph.clear()
                                if translated.strip():
                                    run = paragraph.add_run(translated)
                                    try: run.font.name = font_styles.get('name', None)
                                    except Exception: pass
                                    try: run.font.size = font_styles.get('size', None)
                                    except Exception: pass
                                    try: run.bold = font_styles.get('bold', None)
                                    except Exception: pass
                                    try: run.italic = font_styles.get('italic', None)
                                    except Exception: pass
                                    try: run.underline = font_styles.get('underline', None)
                                    except Exception: pass


        # Save the modified document to a BytesIO object
        output = io.BytesIO()
        doc.save(output)
        output.seek(0) # Rewind the buffer to the beginning
        return output
    except Exception as e:
        print(f"Error translating DOCX: {e}")
        flash(f"Error translating DOCX file: {e}")
        return None

def translate_pptx_in_memory(file_stream, target_lang, model_name, model_type, detected_lang):
    """Translates PPTX file from a BytesIO stream and returns a BytesIO object."""
    try:
        file_stream.seek(0)
        ppt = Presentation(file_stream)
        for slide in ppt.slides:
            for shape in slide.shapes:
                # Handle text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            translated = translate_text(paragraph.text, target_lang, model_name, model_type, detected_lang)
                            if translated is not None and translated.strip():
                                # Simple replacement of text, might lose some formatting nuances per run
                                paragraph.clear() # Clear existing content
                                new_run = paragraph.add_run()
                                new_run.text = translated
                                # Basic font copy (more complex formatting needs deeper pptx manipulation)
                                # if paragraph.runs:
                                #     source_run = paragraph.runs[0]
                                #     try: new_run.font.name = source_run.font.name
                                #     except: pass
                                #     try: new_run.font.size = source_run.font.size
                                #     except: pass

                # Handle tables within shapes
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in shape.table.cells: # Corrected access to cells via shape.table.cells
                            # A cell can contain multiple paragraphs
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    if paragraph.text.strip():
                                        translated = translate_text(paragraph.text, target_lang, model_name, model_type, detected_lang)
                                        if translated is not None and translated.strip():
                                            paragraph.clear()
                                            new_run = paragraph.add_run()
                                            new_run.text = translated

                # Handle grouped shapes (recurse if needed, but simple text frame check is done above)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for subshape in shape.shapes:
                        if subshape.has_text_frame:
                             for paragraph in subshape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    translated = translate_text(paragraph.text, target_lang, model_name, model_type, detected_lang)
                                    if translated is not None and translated.strip():
                                        paragraph.clear()
                                        new_run = paragraph.add_run()
                                        new_run.text = translated


        output = io.BytesIO()
        ppt.save(output)
        output.seek(0)
        return output
    except Exception as e:
        print(f"Error translating PPTX: {e}")
        flash(f"Error translating PPTX file: {e}")
        return None


def translate_excel_in_memory(file_stream, target_lang, model_name, model_type, detected_lang):
    """Translates Excel file from a BytesIO stream and returns a BytesIO object."""
    try:
        file_stream.seek(0)
        df = pd.read_excel(file_stream)

        # Apply translation to each cell that is not NaN
        translated_df = df.applymap(lambda x:
            # Check if cell is not pandas NaN and is not an empty string after stripping
            translate_text(str(x).strip(), target_lang, model_name, model_type, detected_lang)
            if pd.notna(x) and str(x).strip() is not None and str(x).strip() != ""
            else x # Keep original value (including NaN, numbers, dates etc. that weren't strings)
        )

        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            # Setting index=False prevents pandas from writing the DataFrame index as a column
            translated_df.to_excel(writer, index=False)
        output.seek(0)
        return output
    except Exception as e:
        print(f"Error translating Excel: {e}")
        flash(f"Error translating Excel file: {e}")
        return None

# --------------------
# Flask Route Handlers (MODIFIED for GCS processing)
# --------------------

# Use Flask's `g` object to store request-specific data like the request_id
# This is good practice for concurrent requests
@app.before_request
def before_request_func():
    g.request_id = uuid.uuid4().hex # Generate unique ID per request
    g.gcs_temp_paths = [] # List to store paths to clean up (primarily for uploaded files)

@app.after_request
def after_request_func(response):
    # This function runs AFTER the response is sent (or at least prepared).
    # It's used here to clean up the uploaded file(s).
    # The translated file cleanup is handled by GCS Lifecycle Management.
    if gcs_bucket and storage_client and hasattr(g, 'gcs_temp_paths') and g.gcs_temp_paths:
        print(f"Cleanup: Deleting {len(g.gcs_temp_paths)} GCS objects for request {g.request_id} via @after_request")
        blobs_to_delete = [gcs_bucket.blob(path) for path in g.gcs_temp_paths]
        try:
            # batch delete blobs (this is a synchronous operation)
            gcs_bucket.delete_blobs(blobs_to_delete)
            print(f"Cleanup successful for request {g.request_id}")
        except Exception as e_cleanup:
            # Log cleanup errors but don't fail the request
            print(f"GCS Cleanup error for request {g.request_id} (after_request): {e_cleanup}")
    return response


@app.route("/", methods=["GET", "POST"])
def index():
    """Handles file upload, translation, and download link."""
    # Clear flashed messages on GET requests to prevent them persisting across page loads unnecessarily
    if request.method == "GET":
        _ = get_flashed_messages() # Consume messages

    # Check GCS and Gemini availability upfront for POST requests
    if request.method == "POST":
        if not gcs_available:
            flash("Google Cloud Storage is not configured or available. Cannot process files.")
            return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)

        if not gemini_configured or not GEMINI_MODEL:
            flash("Gemini API is not configured. Please set GOOGLE_API_KEY and GEMINI_MODEL in .env.")
            return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)

        file = request.files.get("file")
        target_lang = request.form.get("target_language")

        if not file or file.filename == "":
            flash("No file selected!")
            return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)
        if not target_lang:
             flash("No target language selected!")
             return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)

        safe_filename = secure_filename(file.filename)
        file_extension = os.path.splitext(safe_filename)[1].lower()

        if file_extension not in [".docx", ".pptx", ".xlsx"]:
            flash("Unsupported file type! Only .docx, .pptx, .xlsx are supported.")
            return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)

        # Use the request ID for a temporary GCS folder
        upload_gcs_path = f"{g.request_id}/uploaded/{safe_filename}"
        translated_gcs_path = f"{g.request_id}/translated/translated_{safe_filename}"

        try:
            # --- Step 1: Upload original file to GCS ---
            print(f"Uploading {safe_filename} to gs://{GCS_BUCKET_NAME}/{upload_gcs_path}")
            upload_blob = gcs_bucket.blob(upload_gcs_path)
            # Ensure the file stream is at the beginning before uploading
            file.stream.seek(0)
            upload_blob.upload_from_file(file.stream)
            print(f"Upload complete for {safe_filename}.")
            # Add ONLY the uploaded file path to cleanup list for @after_request
            g.gcs_temp_paths.append(upload_gcs_path)

            # --- Step 2: Read file content from GCS for processing ---
            print(f"Reading content from gs://{GCS_BUCKET_NAME}/{upload_gcs_path} for language detection and translation...")
            file_content_for_detection = ""
            # Read the uploaded blob back into a BytesIO to pass to existing readers
            uploaded_file_stream = blob_to_bytesio(upload_blob)

            if file_extension == ".docx":
                file_content_for_detection = read_text_from_docx(uploaded_file_stream)
            elif file_extension == ".pptx":
                file_content_for_detection = read_text_from_pptx(uploaded_file_stream)
            elif file_extension == ".xlsx":
                file_content_for_detection = read_text_from_excel(uploaded_file_stream)

            # Reset the stream again after reading for detection, before passing to translator
            if uploaded_file_stream:
                 uploaded_file_stream.seek(0)

            # --- Step 3: Language detection ---
            detected_language = None
            if file_content_for_detection and file_content_for_detection.strip():
                detected_language = detect_language(file_content_for_detection)
                if detected_language:
                    flash(f"Detected language: {detected_language.upper()}")
                else:
                    flash("Could not confidently detect language.")
            else:
                flash("Could not extract sufficient text from the file for language detection. Proceeding with translation...")
                detected_language = None

            # --- Step 4: Translate the file using the content read from GCS ---
            print(f"Translating content to {target_lang}...")
            model_type = "google" # Still indicates the API type
            model_name = GEMINI_MODEL
            translated_file_stream = None

            # Pass the stream (read from GCS) to the translation functions
            if file_extension == ".docx":
                translated_file_stream = translate_docx_in_memory(uploaded_file_stream, target_lang, model_name, model_type, detected_language)
            elif file_extension == ".pptx":
                translated_file_stream = translate_pptx_in_memory(uploaded_file_stream, target_lang, model_name, model_type, detected_language)
            elif file_extension == ".xlsx":
                translated_file_stream = translate_excel_in_memory(uploaded_file_stream, target_lang, model_name, model_type, detected_language)

            if translated_file_stream:
                # --- Step 5: Upload the translated file stream to GCS ---
                print(f"Uploading translated file to gs://{GCS_BUCKET_NAME}/{translated_gcs_path}")
                translated_blob = gcs_bucket.blob(translated_gcs_path)
                translated_file_stream.seek(0) # Ensure stream is at the beginning
                translated_blob.upload_from_file(translated_file_stream)
                print(f"Upload complete for translated file.")
                # --- FIX: Do NOT add translated_gcs_path to g.gcs_temp_paths ---
                # It will be cleaned up by GCS Lifecycle Policy configured in the bucket settings.
                # g.gcs_temp_paths.append(translated_gcs_path) # <-- This line is removed/commented out

                # --- Step 6: Store only the GCS path and filename in session ---
                # This replaces storing the large binary data
                file_id = str(uuid.uuid4())
                session[file_id] = {
                    'gcs_path': translated_gcs_path,
                    'filename': f"translated_{safe_filename}"
                }

                # Flash success message only if no errors/blocks occurred during translation
                existing_messages = get_flashed_messages()
                if not any("Error" in msg or "fail" in msg.lower() or "blocked" in msg.lower() for msg in existing_messages):
                    flash("Translation completed successfully! Click the link below to download.")

                # Return the template with the file_id
                return render_template("index.html", languages=LANGUAGES, file_id=file_id, gcs_available=gcs_available, gemini_configured=gemini_configured)
            else:
                # translated_file_stream is None if translation failed
                flash("Translation failed.")
                # Note: Cleanup for uploaded file is handled by @after_request
                return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)

        except Exception as e:
            # Catch any unexpected errors during processing
            print(f"Error during processing (Request ID: {g.request_id}): {e}")
            # Add exception details to flash message
            flash(f"An error occurred during processing: {e}")
            # Note: Cleanup for any files uploaded so far is handled by @after_request
            return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)


    # GET request - render the initial form
    return render_template("index.html", languages=LANGUAGES, gcs_available=gcs_available, gemini_configured=gemini_configured)


@app.route("/download/<file_id>")
def download_file(file_id):
    """Serves translated file downloaded from GCS."""
    # Check GCS availability before attempting download
    if not gcs_available:
        flash("Google Cloud Storage is not configured or available. Cannot download file.")
        return redirect(url_for("index"))

    try:
        # Retrieve GCS path and filename from session using the file_id
        file_info = session.get(file_id)
        if not file_info or 'gcs_path' not in file_info or 'filename' not in file_info:
            flash("File not found or download link expired.")
            return redirect(url_for("index"))

        gcs_path = file_info['gcs_path']
        filename = file_info['filename']

        print(f"Attempting to download gs://{GCS_BUCKET_NAME}/{gcs_path} for request {g.request_id}")

        # --- Step 7: Download the translated file from GCS ---
        translated_blob = gcs_bucket.blob(gcs_path)
        try:
             # Check if the blob exists before trying to download
             if not translated_blob.exists():
                 print(f"Error: Blob not found at gs://{GCS_BUCKET_NAME}/{gcs_path}")
                 flash("Error: Translated file not found (it may have expired).") # Add possibility of expiry
                 session.pop(file_id, None) # Clean up session entry
                 return redirect(url_for("index"))

             # Download blob content into a BytesIO object
             output_stream = io.BytesIO()
             translated_blob.download_to_file(output_stream)
             output_stream.seek(0) # Rewind the buffer

             print(f"Successfully downloaded {gcs_path}. Serving file {filename}.")

             # --- Step 8: Serve the BytesIO stream ---
             response = send_file(
                 output_stream,
                 as_attachment=True,
                 download_name=filename,
                 mimetype='application/octet-stream' # Generic binary type, browser will handle based on extension
             )

             # --- Cleanup: Remove session entry immediately after serving ---
             # GCS cleanup of the blob itself relies on GCS Lifecycle Management.
             # We just clean up the reference in the session.
             session.pop(file_id, None) # Remove the session entry once the download is served

             return response

        except NotFound:
             # This might still happen if GCS deletes it between exists() and download_to_file(), though less likely
             print(f"Error: Blob not found at gs://{GCS_BUCKET_NAME}/{gcs_path} during download.")
             flash("Error: Translated file not found (it may have expired).")
             session.pop(file_id, None)
             return redirect(url_for("index"))
        except Exception as e:
            print(f"Error serving file from GCS {gcs_path} (Request ID: {g.request_id}): {e}")
            flash("Error serving translated file.")
            session.pop(file_id, None) # Clean up session entry on error
            return redirect(url_for("index"))

    except ValueError:
        # Handle invalid file_id format (not a valid UUID)
        flash("Invalid download link.")
        return redirect(url_for("index"))
    except Exception as e:
        # Catch any unexpected errors in the download route
        print(f"Unexpected error in download route (Request ID: {g.request_id}): {e}")
        flash("An unexpected error occurred during download.")
        return redirect(url_for("index"))


if __name__ == "__main__":
    print("\n---------------------------------------------------")
    print("File Translation App Starting...")
    print("---------------------------------------------------")
    print("Configuration Status:")
    if GOOGLE_API_KEY:
        print("- Google Gemini API Key: Found")
        print(f"- Gemini Model: {GEMINI_MODEL or 'Not Set'}")
    else:
        print("- Google Gemini API Key: NOT Found")

    if GCS_BUCKET_NAME:
        print(f"- GCS Bucket Name: {GCS_BUCKET_NAME}")
        if GOOGLE_CLOUD_PROJECT:
             print(f"- Google Cloud Project: {GOOGLE_CLOUD_PROJECT}")
             if gcs_available:
                  print("- GCS Client: Initialized Successfully")
             else:
                  print("- GCS Client: Failed Initialization")
        else:
             print("- Google Cloud Project: NOT Found")
             print("- GCS Client: Not Initialized")
    else:
        print("- GCS Bucket Name: NOT Found")
        print("- GCS Client: Not Initialized")

    if os.getenv("SECRET_KEY"):
        print("- Flask Secret Key: Found")
    else:
        print("- Flask Secret Key: NOT Found (Using random key for this run)")

    print("---------------------------------------------------\n")

    # Consider using a production WSGI server like Gunicorn in Cloud Run
    # For local development, the Flask dev server is fine.
    debug_env = os.environ.get('FLASK_DEBUG', 'False').lower()
    debug_mode = debug_env in ['true', '1', 't', 'yes']

    # In Cloud Run, the PORT environment variable is set
    port = int(os.environ.get('PORT', 5000))
    host = '0.0.0.0' # Listen on all interfaces for cloud deployment

    print(f"Starting Flask app on http://{host}:{port} (Debug: {debug_mode})")

    # Ensure session is configured before running
    # This is done above, but double check if needed based on SESSION_TYPE
    # from flask.sessions import FileSystemSessionInterface
    # app.session_interface = FileSystemSessionInterface(app.config['SESSION_FILE_DIR'], app.config['SECRET_KEY'])


    # When running in production (e.g., Gunicorn), __name__ != '__main__',
    # so this block is skipped and the WSGI server calls app.
    # For local testing, this uses the Flask dev server.
    if __name__ == '__main__':
         # If using filesystem sessions, ensure the directory exists
         # if app.config['SESSION_TYPE'] == 'filesystem':
         #    os.makedirs(app.config['SESSION_FILE_DIR'], exist_ok=True)
         app.run(host=host, port=port, debug=debug_mode)
    # else:
        # If not __main__, it means a WSGI server is running it (like Gunicorn)
        # No need to call app.run() here.
        # Ensure necessary startup logic (like GCS client init) happens outside __main__
        # if it wasn't already, but our current init is global so it's fine.